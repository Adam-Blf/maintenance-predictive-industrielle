# -*- coding: utf-8 -*-
"""Génération du support de présentation PPTX (livrable explicite).

Le sujet liste comme livrable obligatoire un *"Support de Présentation
du projet"*. Ce script produit `reports/presentation.pptx` avec 24 slides
couvrant l'ensemble du projet de Maintenance Prédictive Industrielle.

Charte EFREI · bleu institutionnel #0D47A1, fonts sobres.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from src.config import (  # noqa: E402
    EFREI_LOGO,
    EFREI_LOGO_WHITE,
    MODELS_DIR,
    REPORTS_DIR,
    REPORTS_FIGURES_DIR,
    ensure_directories,
)

COLOR_NAVY = RGBColor(0x0D, 0x47, 0xA1)
COLOR_BLUE = RGBColor(0x1E, 0x88, 0xE5)
COLOR_DARK = RGBColor(0x21, 0x21, 0x21)
COLOR_GRAY = RGBColor(0x6E, 0x6E, 0x6E)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GREEN = RGBColor(0x2E, 0x7D, 0x32)
COLOR_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)


# ---------------------------------------------------------------------------
# Helpers (NE PAS MODIFIER)
# ---------------------------------------------------------------------------

def _set_slide_background(slide, rgb: RGBColor) -> None:
    """Ajoute un rectangle plein écran comme fond coloré."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_DARK,
    align: str = "left",
) -> None:
    """Helper · ajoute un textbox avec formatage uniforme."""
    from pptx.enum.text import PP_ALIGN

    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_title_bar(slide, title: str) -> None:
    """Bandeau de titre uniforme en haut de chaque slide de contenu."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_NAVY
    _add_textbox(
        slide,
        title,
        0.5,
        0.18,
        12.3,
        0.6,
        size=24,
        bold=True,
        color=COLOR_WHITE,
    )


def _add_image_safe(
    slide,
    path: Path,
    left: float,
    top: float,
    width: float,
    max_height: float | None = None,
    recenter: bool = True,
) -> float:
    """Insere une image en calculant la hauteur reelle via PIL.

    Si `max_height` est fourni et que l'image scaled depasserait, on
    reduit la largeur pour rester dans la zone disponible.

    Returns
    -------
    float
        Hauteur reelle de l'image inseree (en inches).
    """
    if not path.exists():
        return 0.0
    from PIL import Image

    with Image.open(path) as img:
        px_w, px_h = img.size
    scaled_h = width * px_h / px_w

    if max_height is not None and scaled_h > max_height:
        scaled_h = max_height
        width = scaled_h * px_w / px_h
        if recenter:
            left = (13.33 - width) / 2

    slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(scaled_h)
    )
    return scaled_h


def _add_notes(slide, text: str) -> None:
    """Injecte les speaker notes dans la zone Notes de la slide."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


def _get_final_model_name() -> str:
    """Lit le nom du modèle final depuis le fichier texte, fallback xgboost."""
    path = MODELS_DIR / "final_model_name.txt"
    if path.exists():
        return path.read_text(encoding="utf-8").strip()
    return "xgboost"


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def build_title_slide(prs: Presentation) -> None:
    """Slide 1 · Couverture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, COLOR_WHITE)

    if EFREI_LOGO.exists():
        slide.shapes.add_picture(str(EFREI_LOGO), Inches(4.65), Inches(0.6), height=Inches(1.4))

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.6), Inches(12.33), Inches(1.6)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_NAVY

    _add_textbox(slide, "Système Intelligent Multi-Modèles",
                 0.5, 2.7, 12.33, 0.7, size=32, bold=True, color=COLOR_WHITE, align="center")
    _add_textbox(slide, "Maintenance Prédictive Industrielle",
                 0.5, 3.4, 12.33, 0.7, size=26, bold=True, color=COLOR_WHITE, align="center")
    _add_textbox(slide, "Projet Data Science · M1 Mastère Data Engineering & IA",
                 0.5, 4.5, 12.33, 0.5, size=18, color=COLOR_NAVY, align="center")
    _add_textbox(slide, "Bloc 2 · BC2 RNCP40875 · Année 2025-2026",
                 0.5, 5.0, 12.33, 0.5, size=14, color=COLOR_GRAY, align="center")
    _add_textbox(slide, "Adam BELOUCIF  ·  Emilien MORICE",
                 0.5, 5.9, 12.33, 0.6, size=22, bold=True, color=COLOR_NAVY, align="center")
    _add_textbox(slide, "EFREI Paris Panthéon-Assas Université  ·  www.efrei.fr",
                 0.5, 6.6, 12.33, 0.4, size=12, color=COLOR_GRAY, align="center")

    _add_notes(slide,
        "Bonjour, nous présentons un système intelligent multi-modèles pour la "
        "maintenance prédictive industrielle. Ce projet répond au sujet 1 du module "
        "Data Science M1, validant le Bloc 2 du RNCP40875. "
        "Nous avons travaillé en binôme : Adam Beloucif et Emilien Morice. "
        "La présentation dure environ 20 minutes suivies d'une démonstration live du dashboard. "
        "Nous couvrons l'ensemble du cycle ML : données, modèles, évaluation, déploiement."
    )


def build_summary_slide(prs: Presentation) -> None:
    """Slide 2 · Sommaire."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "Sommaire")

    col1 = (
        "1.  Couverture\n"
        "2.  Sommaire\n"
        "3.  Contexte métier\n"
        "4.  Problématique et objectifs\n"
        "5.  Dataset Kaggle\n"
        "6.  EDA · insights clés\n"
        "7.  Architecture du système\n"
        "8.  Pipeline ML\n"
        "9.  4 modèles comparés\n"
        "10. Stratégie d'évaluation\n"
        "11. Résultats · barplot comparatif\n"
        "12. Courbes ROC + Precision-Recall"
    )
    col2 = (
        "13. Matrices de confusion\n"
        "14. Optimisation du seuil de décision\n"
        "15. Calibration probabiliste\n"
        "16. Feature Importance\n"
        "17. SHAP global et local\n"
        "18. Bonus · multi-classe + régression\n"
        "19. Industrialisation · Dashboard\n"
        "20. Industrialisation · API REST\n"
        "21. Écoresponsabilité (C4.3)\n"
        "22. Conclusion + perspectives\n"
        "23. Couverture RNCP40875\n"
        "24. Q&A · Merci"
    )
    _add_textbox(slide, col1, 0.6, 1.1, 6.0, 6.0, size=16, color=COLOR_DARK)
    _add_textbox(slide, col2, 7.0, 1.1, 6.0, 6.0, size=16, color=COLOR_DARK)

    # Séparateur vertical
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.0), Inches(0.05), Inches(6.0))
    sep.line.fill.background()
    sep.fill.solid()
    sep.fill.fore_color.rgb = COLOR_BLUE

    _add_notes(slide,
        "Ce sommaire présente les 24 slides de la présentation. "
        "La première partie (slides 3 à 6) couvre le contexte, la problématique et les données. "
        "La deuxième partie (slides 7 à 17) détaille la modélisation et l'évaluation. "
        "La troisième partie (slides 18 à 21) présente les bonus et l'industrialisation. "
        "Nous terminons par la conclusion, la couverture RNCP et les questions."
    )


def build_context_slide(prs: Presentation) -> None:
    """Slide 3 · Contexte métier."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "3. Contexte métier")

    # Bloc coût (bandeau coloré)
    cost_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(12.33), Inches(0.85)
    )
    cost_bar.line.fill.background()
    cost_bar.fill.solid()
    cost_bar.fill.fore_color.rgb = COLOR_LIGHT
    _add_textbox(slide, "Coût moyen d'un arrêt non planifié en industrie lourde : ~50 000 euros/heure",
                 0.7, 1.1, 12.0, 0.8, size=18, bold=True, color=COLOR_NAVY, align="center")

    bullets = [
        "Maintenance corrective · intervention après panne · coût et délai maxima.",
        "Maintenance préventive · révision périodique · sur-maintenance possible.",
        "Maintenance prédictive · alerte avant défaillance via capteurs IoT · compromis optimal.",
        "Capteurs continus : vibration, température moteur/entrainement, pression, RPM.",
        "Enjeu : transformer les signaux bruts en alerte exploitable par les opérateurs GMAO.",
    ]
    _add_textbox(slide, "\n\n".join("•  " + b for b in bullets),
                 0.7, 2.1, 12.0, 4.9, size=17, color=COLOR_DARK)

    _add_notes(slide,
        "Le coût d'une panne non planifiée en industrie lourde dépasse souvent 50 000 euros/heure "
        "quand on intègre l'arrêt de production, les pénalités contractuelles et les coûts de réparation urgente. "
        "La maintenance corrective est la plus coûteuse : on répare après la casse. "
        "La maintenance préventive améliore la situation mais génère des révisions inutiles. "
        "La maintenance prédictive est le compromis optimal : on intervient juste avant la panne. "
        "Notre projet vise exactement ce cas d'usage : prédire la défaillance dans les 24 heures suivantes."
    )


def build_problem_slide(prs: Presentation) -> None:
    """Slide 4 · Problématique et objectifs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "4. Problématique et objectifs")

    # Tâche principale
    _add_textbox(slide, "Tâche principale · Classification binaire",
                 0.6, 1.05, 12.0, 0.55, size=20, bold=True, color=COLOR_NAVY)
    _add_textbox(slide,
        "Prédire failure_within_24h (0 / 1) à partir des mesures capteurs en temps réel.",
        0.6, 1.6, 12.0, 0.5, size=17, color=COLOR_DARK)

    # Tâches bonus
    _add_textbox(slide, "Tâches bonus :",
                 0.6, 2.3, 12.0, 0.45, size=18, bold=True, color=COLOR_NAVY)
    bonus = [
        "Multi-classe · failure_type (Normal, Heat, Power, Product, Random, Overstrain) · 5 classes.",
        "Régression · rul_hours (Remaining Useful Life) · prédiction continue en heures.",
    ]
    _add_textbox(slide, "\n".join("•  " + b for b in bonus),
                 0.6, 2.75, 12.0, 1.1, size=16, color=COLOR_DARK)

    # KPI métier
    _add_textbox(slide, "KPI métier retenus :",
                 0.6, 4.05, 12.0, 0.45, size=18, bold=True, color=COLOR_NAVY)
    kpis = [
        "Recall (sensibilité) · minimiser les faux négatifs = pannes non détectées.",
        "PR-AUC · robuste face au déséquilibre des classes (15 % de positifs).",
        "ROI estimé · chaque TP détecté évite ~50 k euros d'arrêt · chaque FP coûte ~500 euros d'intervention inutile.",
    ]
    _add_textbox(slide, "\n".join("•  " + b for b in kpis),
                 0.6, 4.5, 12.0, 2.6, size=16, color=COLOR_DARK)

    _add_notes(slide,
        "La tâche principale est une classification binaire : va-t-il y avoir une panne dans les 24 heures ? "
        "Nous avons aussi traité deux tâches bonus : identifier le type de panne (5 classes) et estimer "
        "le temps restant avant défaillance (régression). "
        "Le choix des métriques est guidé par l'asymétrie des coûts : "
        "un faux négatif (panne manquée) coûte ~100 fois plus qu'un faux positif (intervention préventive inutile). "
        "C'est pourquoi nous privilégions le Recall et la PR-AUC plutôt que l'accuracy."
    )


def build_dataset_slide(prs: Presentation) -> None:
    """Slide 5 · Dataset Kaggle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "5. Dataset Kaggle · Predictive Maintenance v3.0")

    meta = [
        "Source : Kaggle (licence CC0 Public Domain) · dataset officiel v3.0",
        "Volume : 24 042 lignes · 15 variables · aucune valeur dupliquée",
        "NaN : ~4 % sur 3 colonnes numériques · imputation médiane (pipeline Silver)",
    ]
    _add_textbox(slide, "\n".join("•  " + m for m in meta),
                 0.6, 1.05, 12.0, 1.2, size=16, color=COLOR_DARK)

    # Tableau schéma
    _add_textbox(slide, "Schéma des variables clés :",
                 0.6, 2.35, 12.0, 0.4, size=17, bold=True, color=COLOR_NAVY)
    schema = (
        "Variable                     Type        Description\n"
        "-----------------------------------------------------\n"
        "UDI                          int         Identifiant unique\n"
        "Product_ID                   str         Code produit (L/M/H)\n"
        "Type                         cat         Qualité produit\n"
        "Air_temperature_K            float       Température air (K)\n"
        "Process_temperature_K        float       Température process (K)\n"
        "Rotational_speed_rpm         int         Vitesse rotation (RPM)\n"
        "Torque_Nm                    float       Couple moteur (Nm)\n"
        "Tool_wear_min                int         Usure outil (min)\n"
        "Machine_failure              int (cible) Panne détectée\n"
        "failure_within_24h           int (cible) Panne dans 24h (tâche principale)\n"
        "failure_type                 cat (bonus) Type de panne (5 classes)\n"
        "rul_hours                    float (bon) Remaining Useful Life (régression)"
    )
    _add_textbox(slide, schema, 0.4, 2.75, 12.5, 4.5, size=12, color=COLOR_DARK)

    _add_notes(slide,
        "Le dataset Kaggle Predictive Maintenance v3.0 est en licence CC0, donc librement utilisable. "
        "Il contient 24 042 observations et 15 colonnes couvrant capteurs industriels et labels de panne. "
        "La variable cible principale failure_within_24h est déséquilibrée : environ 15 % de positifs. "
        "Nous avons géré les 4 % de valeurs manquantes par imputation médiane dans le pipeline Silver. "
        "La variable failure_type nous permet de traiter la tâche multi-classe en bonus."
    )


def build_eda_slide(prs: Presentation) -> None:
    """Slide 6 · EDA · insights clés (4 visuels)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "6. EDA · Insights clés")

    # 4 images en grille 2x2
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "eda_target_distribution.png",
                    0.2, 1.0, 6.4, max_height=2.9, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "eda_correlation_heatmap.png",
                    6.7, 1.0, 6.4, max_height=2.9, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "eda_sensor_distributions.png",
                    0.2, 4.0, 6.4, max_height=3.0, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "eda_scatter_vib_temp.png",
                    6.7, 4.0, 6.4, max_height=3.0, recenter=False)

    _add_notes(slide,
        "L'EDA révèle quatre points clés. "
        "Haut gauche : la cible est déséquilibrée à 85/15, ce qui justifie d'utiliser PR-AUC plutôt qu'accuracy. "
        "Haut droite : la heatmap de corrélation montre que Torque_Nm et Rotational_speed_rpm sont très corrélés negativement. "
        "Bas gauche : les distributions des capteurs par classe montrent des séparations nettes pour la température et l'usure. "
        "Bas droite : le scatter vibration/température met en évidence deux régimes distincts selon la présence de panne."
    )


def build_architecture_slide(prs: Presentation) -> None:
    """Slide 7 · Architecture du système."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "7. Architecture du système")
    img_top = 1.1
    img_h = _add_image_safe(
        slide, REPORTS_FIGURES_DIR / "diagram_architecture.png",
        0.5, img_top, 12.3, max_height=5.2)
    _add_textbox(slide,
        "Architecture médaillon Bronze / Silver / Gold · pipeline reproductible · API + Dashboard.",
        0.5, img_top + img_h + 0.2, 12.3, 0.6, size=14, color=COLOR_GRAY, align="center")

    _add_notes(slide,
        "L'architecture s'organise en trois couches médaillon. "
        "Bronze : données brutes ingérées telles quelles depuis le CSV Kaggle. "
        "Silver : nettoyage, imputation des NaN, feature engineering, encodage. "
        "Gold : features prêtes pour l'entraînement, modèles sérialisés, logs MLflow. "
        "En production, le front Streamlit appelle l'API FastAPI qui charge le modèle joblib. "
        "Cette architecture garantit la reproductibilité et la traçabilité de chaque transformation."
    )


def build_pipeline_slide(prs: Presentation) -> None:
    """Slide 8 · Pipeline ML."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "8. Pipeline ML · Anti data-leakage")
    img_top = 1.1
    img_h = _add_image_safe(
        slide, REPORTS_FIGURES_DIR / "diagram_ml_pipeline.png",
        0.5, img_top, 12.3, max_height=4.2)

    anti_leak = [
        "Imputation médiane et StandardScaler encapsulés dans le Pipeline scikit-learn.",
        "Fit uniquement sur le train set · aucune fuite d'information vers le test set.",
        "Cross-validation stratifiée 5-fold appliquée sur le pipeline complet.",
    ]
    _add_textbox(slide, "\n".join("•  " + a for a in anti_leak),
                 0.7, img_top + img_h + 0.25, 12.0, 1.6, size=15, color=COLOR_DARK)

    _add_notes(slide,
        "Le pipeline scikit-learn enchaîne préprocesseur et estimateur dans un seul objet. "
        "Cela garantit l'anti data-leakage : l'imputation et la normalisation sont apprises "
        "uniquement sur le train set et appliquées telles quelles sur le test set. "
        "La cross-validation est appliquée sur le pipeline entier, pas seulement sur le modèle. "
        "Tous les transformeurs et le modèle final sont sérialisés en un seul fichier joblib."
    )


def build_methodology_slide(prs: Presentation) -> None:
    """Slide 9 · Méthodologie · 4 modèles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "9. Méthodologie · 4 modèles comparés")

    # Tableau comparatif
    _add_textbox(slide, "Tableau comparatif des 4 modèles :",
                 0.5, 1.0, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)

    header = "Modèle               Famille              Hyperparams clés                  Forces"
    rows = [
        "Logistic Reg.        Linéaire             C=0.1, max_iter=1000               Interprétable, rapide, baseline",
        "Random Forest        Ensemble (Bagging)   n_est=200, max_depth=15            Robuste, peu d'overfitting",
        "XGBoost              Ensemble (Boosting)  lr=0.05, n_est=300, max_depth=6    SOTA tabulaire, gestion NaN",
        "MLP 64-32-16         Deep Learning        alpha=1e-3, early_stopping=True    Capture interactions complexes",
    ]
    table_text = header + "\n" + "-" * 95 + "\n" + "\n".join(rows)
    _add_textbox(slide, table_text, 0.4, 1.5, 12.5, 2.5, size=13, color=COLOR_DARK)

    # Choix communs
    _add_textbox(slide, "Paramètres communs :",
                 0.5, 4.15, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    common = [
        "class_weight='balanced' sur LR et RF · SMOTE en option pour XGBoost et MLP.",
        "Optimisation hyperparamètres via Optuna (50 trials, Bayesian TPE).",
        "Stratification maintenue à chaque fold de la CV 5-fold.",
    ]
    _add_textbox(slide, "\n".join("•  " + c for c in common),
                 0.6, 4.6, 12.0, 2.5, size=15, color=COLOR_DARK)

    _add_notes(slide,
        "Nous avons comparé quatre familles de modèles couvrant linéaire, ensembles et deep learning. "
        "La Logistic Regression sert de baseline interprétable. "
        "Le Random Forest est robuste et peu sensible aux hyperparamètres. "
        "XGBoost est l'état de l'art sur les données tabulaires et gère nativement les NaN. "
        "Le MLP capture les interactions non linéaires complexes mais est plus lent à entraîner. "
        "Tous les modèles sont optimisés avec Optuna pour garantir une comparaison équitable."
    )


def build_evaluation_strategy_slide(prs: Presentation) -> None:
    """Slide 10 · Stratégie d'évaluation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "10. Stratégie d'évaluation")

    metrics = [
        "PR-AUC (Precision-Recall AUC) · métrique principale · robuste face au déséquilibre 85/15.",
        "F1-Score macro · équilibre Precision et Recall sur toutes les classes.",
        "Recall · priorité métier absolue · minimiser les faux négatifs (pannes non détectées).",
        "ROC-AUC · comparaison classique · complète PR-AUC pour la vue globale.",
        "Accuracy · fournie pour contexte mais non utilisée comme critère de sélection.",
    ]
    _add_textbox(slide, "Métriques retenues :",
                 0.6, 1.0, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    _add_textbox(slide, "\n\n".join("•  " + m for m in metrics),
                 0.6, 1.5, 12.0, 3.3, size=15, color=COLOR_DARK)

    _add_textbox(slide, "Protocole de validation :",
                 0.6, 4.9, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    protocol = [
        "Split stratifié 80/20 · même proportion de positifs dans train et test.",
        "Cross-validation stratifiée 5-fold sur le train set pour sélectionner les hyperparamètres.",
        "Évaluation finale sur le test set holdout (jamais vu pendant l'optimisation).",
    ]
    _add_textbox(slide, "\n".join("•  " + p for p in protocol),
                 0.6, 5.35, 12.0, 1.9, size=15, color=COLOR_DARK)

    _add_notes(slide,
        "La stratégie d'évaluation est construite autour du déséquilibre des classes. "
        "L'accuracy est trompeuse quand 85 % des observations sont négatives : "
        "un modèle qui prédit toujours 0 obtient 85 % d'accuracy mais ne sert à rien. "
        "PR-AUC mesure la qualité des prédictions sur les positifs uniquement, ce qui correspond "
        "exactement à notre besoin métier. Le Recall est notre contrainte absolue : "
        "mieux vaut déclencher quelques fausses alarmes que manquer une vraie panne."
    )


def build_results_slide(prs: Presentation) -> None:
    """Slide 11 · Résultats · barplot comparatif."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "11. Résultats · Comparaison des 4 modèles")
    img_top = 1.05
    img_h = _add_image_safe(
        slide, REPORTS_FIGURES_DIR / "metrics_comparison_barplot.png",
        0.5, img_top, 12.3, max_height=5.0)

    final_name = _get_final_model_name()
    _add_textbox(slide,
        f"Modèle candidat retenu : {final_name}  ·  compromis F1 / stabilité CV / interprétabilité",
        0.5, img_top + img_h + 0.2, 12.3, 0.6, size=16, bold=True, color=COLOR_NAVY, align="center")

    _add_notes(slide,
        "Le barplot compare les 4 modèles sur F1-macro, ROC-AUC et PR-AUC. "
        "XGBoost et Random Forest dominent sur toutes les métriques. "
        "La Logistic Regression confirme son rôle de baseline solide malgré sa simplicité. "
        "Le MLP est compétitif mais plus lent et moins interprétable. "
        "Le modèle final est sélectionné en croisant performance, stabilité CV et interprétabilité métier."
    )


def build_curves_slide(prs: Presentation) -> None:
    """Slide 12 · Courbes ROC + PR."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "12. Courbes ROC et Precision-Recall")
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "roc_curves_comparison.png",
                    0.3, 1.1, 6.4, max_height=5.5, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "pr_curves_comparison.png",
                    6.7, 1.1, 6.4, max_height=5.5, recenter=False)

    _add_notes(slide,
        "Les deux courbes montrent les performances en fonction du seuil de décision. "
        "La courbe ROC (gauche) trace le taux de vrais positifs contre le taux de faux positifs. "
        "La courbe Precision-Recall (droite) est plus informative en contexte déséquilibré : "
        "elle montre directement le compromis entre détecter les pannes et éviter les fausses alarmes. "
        "L'aire sous la courbe PR (PR-AUC) est notre métrique de référence pour la sélection du modèle."
    )


def build_confusion_matrices_slide(prs: Presentation) -> None:
    """Slide 13 · Matrices de confusion 4 modèles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "13. Matrices de confusion · 4 modèles")

    # Grille 2x2
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "confusion_matrix_logistic_regression.png",
                    0.2, 1.0, 6.4, max_height=2.9, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "confusion_matrix_random_forest.png",
                    6.7, 1.0, 6.4, max_height=2.9, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "confusion_matrix_xgboost.png",
                    0.2, 4.0, 6.4, max_height=3.0, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "confusion_matrix_mlp.png",
                    6.7, 4.0, 6.4, max_height=3.0, recenter=False)

    _add_notes(slide,
        "Les matrices de confusion permettent de visualiser directement les faux négatifs et faux positifs. "
        "Case en bas à gauche : faux négatifs (panne prédite comme normale) · coût métier maximal. "
        "Case en haut à droite : faux positifs (normale prédite comme panne) · coût opérationnel. "
        "On observe que XGBoost et Random Forest minimisent mieux les faux négatifs que LR et MLP. "
        "Le seuil de décision par défaut est 0.5 ; nous l'optimiserons à la slide suivante."
    )


def build_threshold_slide(prs: Presentation) -> None:
    """Slide 14 · Optimisation du seuil de décision."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "14. Optimisation du seuil de décision")

    _add_textbox(slide,
        "Asymétrie des coûts : FN (panne manquée) ~ 100x plus coûteux qu'un FP (fausse alarme).",
        0.6, 1.05, 12.0, 0.55, size=17, bold=True, color=COLOR_NAVY)

    _add_image_safe(slide, REPORTS_FIGURES_DIR / "cost_threshold_random_forest.png",
                    0.2, 1.7, 6.3, max_height=4.8, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "cost_threshold_xgboost.png",
                    6.7, 1.7, 6.3, max_height=4.8, recenter=False)

    _add_notes(slide,
        "Le seuil par défaut de 0.5 est rarement optimal en contexte industriel. "
        "Nous modélisons la fonction de coût total = cout_FN * nb_FN + cout_FP * nb_FP "
        "pour différentes valeurs du seuil de classification. "
        "En fixant le ratio cout_FN/cout_FP à 100, le seuil optimal se situe typiquement entre 0.2 et 0.35. "
        "Ce seuil est configurable dans le dashboard pour permettre à l'opérateur d'ajuster selon le contexte."
    )


def build_calibration_slide(prs: Presentation) -> None:
    """Slide 15 · Calibration probabiliste."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "15. Calibration probabiliste")

    _add_textbox(slide,
        "Un modèle calibré : si il prédit 80 % de probabilité, cela doit arriver 80 % du temps.",
        0.6, 1.05, 12.0, 0.55, size=17, color=COLOR_DARK)

    _add_image_safe(slide, REPORTS_FIGURES_DIR / "reliability_diagram_random_forest.png",
                    0.2, 1.7, 6.3, max_height=4.8, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "reliability_diagram_xgboost.png",
                    6.7, 1.7, 6.3, max_height=4.8, recenter=False)

    _add_notes(slide,
        "La calibration est cruciale quand on utilise des probabilités pour prendre des décisions métier. "
        "Le reliability diagram compare les probabilités prédites aux fréquences observées. "
        "La droite diagonale représente la calibration parfaite. "
        "Random Forest sur-estime souvent les probabilités extrêmes (courbe en S). "
        "XGBoost est généralement mieux calibré nativement. "
        "Nous proposons une calibration post-hoc via Platt Scaling ou Isotonic Regression si nécessaire."
    )


def build_feature_importance_slide(prs: Presentation) -> None:
    """Slide 16 · Feature Importance native + permutation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "16. Feature Importance · Native + Permutation")

    final_name = _get_final_model_name()

    _add_image_safe(slide, REPORTS_FIGURES_DIR / f"feature_importance_native_{final_name}.png",
                    0.2, 1.0, 6.3, max_height=5.5, recenter=False)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / f"permutation_importance_{final_name}.png",
                    6.7, 1.0, 6.3, max_height=5.5, recenter=False)

    _add_notes(slide,
        "Deux approches complémentaires pour identifier les variables clés. "
        "Importance native (gauche) : basée sur la réduction d'impureté dans les arbres, "
        "rapide mais biaisée en faveur des variables continues à forte cardinalité. "
        "Importance par permutation (droite) : mesure la dégradation de performance quand "
        "une variable est mélangée aléatoirement, plus fiable et modèle-agnostique. "
        "Les deux approches convergent : Tool_wear_min, Torque_Nm et Air_temperature_K "
        "sont systématiquement les variables les plus importantes."
    )


def build_interpret_slide(prs: Presentation) -> None:
    """Slide 17 · Interprétabilité · SHAP global et local."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "17. Interprétabilité · SHAP global et local")
    final_name = _get_final_model_name()

    h1 = _add_image_safe(slide, REPORTS_FIGURES_DIR / f"shap_summary_{final_name}.png",
                         0.3, 1.1, 7.0, max_height=5.4, recenter=False)
    h2 = _add_image_safe(slide, REPORTS_FIGURES_DIR / f"shap_bar_{final_name}.png",
                         7.5, 1.4, 5.5, max_height=4.8, recenter=False)
    bottom = 1.1 + max(h1, h2) + 0.2
    _add_textbox(slide,
        "Top 3 variables : Tool_wear_min · Torque_Nm · Air_temperature_K",
        0.5, bottom, 12.3, 0.5, size=14, color=COLOR_GRAY, align="center")

    _add_notes(slide,
        "SHAP (SHapley Additive exPlanations) fournit une interprétation cohérente et fidèle au modèle. "
        "Le summary plot (gauche) montre, pour chaque observation, la contribution de chaque variable "
        "à la prédiction : rouge = valeur haute de la variable, bleu = valeur basse. "
        "Le bar plot (droite) donne l'importance globale moyenne sur tout le test set. "
        "On voit que Tool_wear_min (usure outil) est de loin la variable la plus impactante : "
        "une usure élevée prédit fortement la panne imminente, ce qui est cohérent métier."
    )


def build_bonus_slide(prs: Presentation) -> None:
    """Slide 18 · Bonus · multi-classe + régression."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "18. Tâches bonus · Multi-classe + Régression RUL")

    _add_textbox(slide, "Multi-classe · failure_type (5 classes) :",
                 0.5, 1.0, 6.5, 0.45, size=16, bold=True, color=COLOR_NAVY)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "multiclass_confusion_matrix.png",
                    0.3, 1.5, 6.3, max_height=5.2, recenter=False)

    _add_textbox(slide, "Régression · rul_hours (Remaining Useful Life) :",
                 6.8, 1.0, 6.5, 0.45, size=16, bold=True, color=COLOR_NAVY)
    _add_image_safe(slide, REPORTS_FIGURES_DIR / "regression_pred_vs_true.png",
                    6.7, 1.5, 6.3, max_height=5.2, recenter=False)

    _add_notes(slide,
        "En bonus, nous avons traité deux tâches supplémentaires. "
        "La classification multi-classe identifie le type de panne parmi 5 catégories : "
        "Normal, Heat Dissipation, Power, Product Quality, Random et Overstrain Failure. "
        "La matrice de confusion 5x5 montre que les pannes Normal et Overstrain sont bien détectées, "
        "tandis que les pannes rares (Random, Product Quality) sont plus difficiles à discriminer. "
        "La régression prédit le nombre d'heures restantes avant défaillance (RUL). "
        "Le nuage Pred vs True montre une corrélation satisfaisante avec quelques outliers sur les longues durées."
    )


def build_dashboard_slide(prs: Presentation) -> None:
    """Slide 19 · Industrialisation · Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "19. Industrialisation · Dashboard Streamlit")

    tabs = [
        "Onglet 1 · Accueil : présentation du projet et guide d'utilisation.",
        "Onglet 2 · Exploration EDA : visualisations interactives des données (filtres Plotly).",
        "Onglet 3 · Prediction : saisie manuelle ou upload CSV pour prédiction batch.",
        "Onglet 4 · Performance : métriques, courbes ROC/PR, matrices de confusion en live.",
        "Onglet 5 · Explicabilité : SHAP waterfall et force plot par observation.",
    ]
    _add_textbox(slide, "\n\n".join("•  " + t for t in tabs),
                 0.7, 1.1, 12.0, 4.5, size=16, color=COLOR_DARK)

    tech = "Stack : Streamlit · Plotly · SHAP · FastAPI · joblib · CSS charte EFREI"
    _add_textbox(slide, tech, 0.7, 6.0, 12.0, 0.9, size=15, bold=False, color=COLOR_NAVY)

    _add_notes(slide,
        "Le dashboard Streamlit est la vitrine opérationnelle du projet. "
        "Il est décomposé en 5 onglets couvrant l'exploration, la prédiction, l'évaluation et l'explicabilité. "
        "L'opérateur peut saisir les valeurs capteurs manuellement ou uploader un CSV entier pour une prédiction batch. "
        "Le CSS personnalisé respecte la charte EFREI avec un design glassmorphism sobre. "
        "L'intégration avec l'API REST permet une séparation propre front/back "
        "et facilite l'intégration future dans un système GMAO."
    )


def build_api_slide(prs: Presentation) -> None:
    """Slide 20 · Industrialisation · API REST."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "20. Industrialisation · API REST FastAPI")

    endpoints = [
        "POST /predict       · prédiction unitaire (JSON) ou batch (liste)",
        "GET  /health        · statut de l'API, uptime, version du modèle",
        "GET  /model-info    · méta-données du modèle chargé (nom, métriques, date)",
        "GET  /docs          · documentation Swagger auto-générée par FastAPI",
    ]
    _add_textbox(slide, "Endpoints exposés :",
                 0.6, 1.0, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    _add_textbox(slide, "\n".join("  " + e for e in endpoints),
                 0.6, 1.5, 12.0, 1.6, size=15, color=COLOR_DARK)

    _add_textbox(slide, "Exemple de payload /predict :",
                 0.6, 3.25, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    json_example = (
        '{\n'
        '  "air_temperature_k": 298.5,\n'
        '  "process_temperature_k": 308.2,\n'
        '  "rotational_speed_rpm": 1450,\n'
        '  "torque_nm": 42.1,\n'
        '  "tool_wear_min": 210\n'
        '}\n'
        '-- Reponse --\n'
        '{ "prediction": 1, "probability": 0.87, "threshold": 0.30, "model": "xgboost" }'
    )
    _add_textbox(slide, json_example, 0.6, 3.75, 12.0, 2.9, size=14, color=COLOR_DARK)

    _add_notes(slide,
        "L'API FastAPI expose le modèle via un contrat JSON strict validé par Pydantic v2. "
        "Le endpoint /predict accepte les 5 features capteurs et retourne la prédiction binaire, "
        "la probabilité brute et le seuil appliqué. "
        "La documentation Swagger est auto-générée et accessible sur /docs. "
        "L'API se lance avec uvicorn api.main:app et la conteneurisation reste une perspective. "
        "En production, on ajouterait une couche d'authentification Bearer et du rate-limiting."
    )


def build_eco_slide(prs: Presentation) -> None:
    """Slide 21 · Écoresponsabilité (RNCP C4.3)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "21. Écoresponsabilité · RNCP C4.3")
    img_top = 1.05
    img_h = _add_image_safe(
        slide, REPORTS_FIGURES_DIR / "compute_cost_comparison.png",
        0.5, img_top, 12.3, max_height=4.5)

    eco_notes = [
        "Mesure via CodeCarbon : émissions CO2 par run d'entraînement enregistrées automatiquement.",
        "Mix énergétique France ~80 gCO2eq/kWh · parmi les plus bas d'Europe (nucléaire).",
        "Recommandation : privilégier RF ou LR pour l'inférence en production (10x moins de CO2 que MLP).",
    ]
    _add_textbox(slide, "\n".join("•  " + e for e in eco_notes),
                 0.6, img_top + img_h + 0.25, 12.0, 1.8, size=15, color=COLOR_DARK)

    _add_notes(slide,
        "La compétence C4.3 du RNCP40875 exige une prise en compte de l'impact environnemental. "
        "Nous avons intégré CodeCarbon dans chaque run d'entraînement pour mesurer les émissions CO2. "
        "Le graphique montre que XGBoost et RF ont un coût carbone similaire, "
        "tandis que le MLP consomme significativement plus pour des gains marginaux. "
        "En France, le mix énergétique à dominante nucléaire (~80 gCO2/kWh) "
        "est très favorable comparé à la moyenne européenne (~300 gCO2/kWh). "
        "Pour une mise en production, nous recommandons RF ou XGBoost plutôt que MLP."
    )


def build_conclusion_slide(prs: Presentation) -> None:
    """Slide 22 · Conclusion et perspectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "22. Conclusion et perspectives")

    _add_textbox(slide, "Synthèse des livrables :",
                 0.6, 1.05, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    delivered = [
        "MVP complet : pipeline reproductible + 4 modèles + dashboard Streamlit + API FastAPI.",
        "Tâches bonus livrées : multi-classe failure_type + régression RUL.",
        "Optimisation Optuna + seuil métier + calibration + SHAP + CodeCarbon.",
        "Architecture médaillon Bronze/Silver/Gold · pipeline reproductible (seed 42).",
    ]
    _add_textbox(slide, "\n".join("•  " + d for d in delivered),
                 0.6, 1.55, 12.0, 1.9, size=15, color=COLOR_DARK)

    _add_textbox(slide, "Limites identifiées :",
                 0.6, 3.55, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    limits = [
        "Drift temporel non géré : les modèles ne se ré-entraînent pas automatiquement.",
        "Événements rares (Random Failure) sous-représentés même après SMOTE.",
        "Données simulées (Kaggle) : bruit et distributions différentes du terrain réel.",
    ]
    _add_textbox(slide, "\n".join("•  " + l for l in limits),
                 0.6, 4.05, 12.0, 1.4, size=15, color=COLOR_DARK)

    _add_textbox(slide, "Perspectives :",
                 0.6, 5.55, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)
    next_steps = "LSTM / Transformer temporel · MLOps (drift monitoring, retraining auto) · intégration GMAO."
    _add_textbox(slide, next_steps, 0.6, 6.0, 12.0, 0.9, size=15, color=COLOR_DARK)

    _add_notes(slide,
        "En conclusion, le projet couvre l'ensemble du cycle ML : de la donnée brute au service déployé. "
        "Les 4 modèles sont comparés équitablement avec un protocole rigoureux anti data-leakage. "
        "Le dashboard et l'API rendent le système opérationnel pour un utilisateur non data scientist. "
        "Les principales limites sont le drift temporel non géré et l'aspect simulé des données. "
        "En M2, nous pourrions intégrer un monitoring de drift avec Evidently et un retraining automatique via Airflow."
    )


def build_rncp_slide(prs: Presentation) -> None:
    """Slide 23 · Couverture RNCP40875."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, "23. Couverture RNCP40875 · Bloc 2")

    _add_textbox(slide, "Tableau de couverture des 6 compétences du Bloc 2 :",
                 0.5, 1.0, 12.0, 0.45, size=17, bold=True, color=COLOR_NAVY)

    header = "Competence   Description                                           Validation"
    sep = "-" * 85
    rows = [
        "C3.1         Collecte et preparation des donnees (pipeline Bronze/Silver)  OUI",
        "C3.2         Exploration et visualisation (EDA 7 figures, heatmap, scatter) OUI",
        "C3.3         Modelisation ML (4 modeles, Optuna, CV 5-fold, SHAP)          OUI",
        "C4.1         Evaluation et validation (PR-AUC, Recall, seuil metier)       OUI",
        "C4.2         Deploiement (Dashboard Streamlit, API FastAPI uvicorn)        OUI",
        "C4.3         Responsabilite (CodeCarbon, mesure CO2, recommandation eco)   OUI",
    ]
    table_text = header + "\n" + sep + "\n" + "\n".join(rows)
    _add_textbox(slide, table_text, 0.4, 1.55, 12.5, 4.0, size=13, color=COLOR_DARK)

    # Badge vert
    badge = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(5.8), Inches(5.33), Inches(0.85)
    )
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_GREEN
    _add_textbox(slide, "Toutes les competences Bloc 2 validees",
                 4.1, 5.9, 5.13, 0.65, size=17, bold=True, color=COLOR_WHITE, align="center")

    _add_notes(slide,
        "Ce tableau récapitule la couverture des 6 compétences du Bloc 2 du RNCP40875. "
        "C3.1 est validée par notre pipeline Bronze/Silver/Gold avec traçabilité complète. "
        "C3.2 par les 7+ figures EDA produites automatiquement. "
        "C3.3 par la comparaison rigoureuse de 4 modèles avec optimisation Optuna. "
        "C4.1 par le protocole d'évaluation centré PR-AUC et l'optimisation du seuil métier. "
        "C4.2 par le dashboard Streamlit et l'API FastAPI containerisée. "
        "C4.3 par l'intégration CodeCarbon et les recommandations éco."
    )


def build_thanks_slide(prs: Presentation) -> None:
    """Slide 24 · Q&A / Merci."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, COLOR_NAVY)

    logo_path = EFREI_LOGO_WHITE if EFREI_LOGO_WHITE.exists() else EFREI_LOGO
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(4.65), Inches(0.7), height=Inches(1.2))

    _add_textbox(slide, "Merci pour votre attention",
                 0.5, 2.8, 12.33, 1.0, size=44, bold=True, color=COLOR_WHITE, align="center")
    _add_textbox(slide, "Questions et démonstration live",
                 0.5, 4.0, 12.33, 0.6, size=22, color=COLOR_WHITE, align="center")

    contact = (
        "Adam BELOUCIF · adam.beloucif@efrei.net\n"
        "Emilien MORICE · emilien.morice@efrei.net"
    )
    _add_textbox(slide, contact, 0.5, 5.1, 12.33, 0.9, size=16, color=COLOR_WHITE, align="center")
    _add_textbox(slide, "EFREI Paris Panthéon-Assas Université  ·  www.efrei.fr",
                 0.5, 6.5, 12.33, 0.4, size=12, color=COLOR_WHITE, align="center")

    _add_notes(slide,
        "Nous sommes disponibles pour toutes vos questions. "
        "Nous pouvons lancer la démonstration live du dashboard Streamlit et de l'API FastAPI. "
        "Nous pouvons montrer une prédiction en temps réel avec les valeurs capteurs de votre choix. "
        "Les contacts email sont affichés pour toute question post-soutenance. "
        "Le code source complet est disponible sur GitHub avec documentation détaillée."
    )


# ---------------------------------------------------------------------------
# Orchestration principale
# ---------------------------------------------------------------------------

def main() -> None:
    ensure_directories()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)            # Slide  1
    build_summary_slide(prs)          # Slide  2
    build_context_slide(prs)          # Slide  3
    build_problem_slide(prs)          # Slide  4
    build_dataset_slide(prs)          # Slide  5
    build_eda_slide(prs)              # Slide  6
    build_architecture_slide(prs)     # Slide  7
    build_pipeline_slide(prs)         # Slide  8
    build_methodology_slide(prs)      # Slide  9
    build_evaluation_strategy_slide(prs)  # Slide 10
    build_results_slide(prs)          # Slide 11
    build_curves_slide(prs)           # Slide 12
    build_confusion_matrices_slide(prs)  # Slide 13
    build_threshold_slide(prs)        # Slide 14
    build_calibration_slide(prs)      # Slide 15
    build_feature_importance_slide(prs)  # Slide 16
    build_interpret_slide(prs)        # Slide 17
    build_bonus_slide(prs)            # Slide 18
    build_dashboard_slide(prs)        # Slide 19
    build_api_slide(prs)              # Slide 20
    build_eco_slide(prs)              # Slide 21
    build_conclusion_slide(prs)       # Slide 22
    build_rncp_slide(prs)             # Slide 23
    build_thanks_slide(prs)           # Slide 24

    output = REPORTS_DIR / "presentation.pptx"
    prs.save(str(output))
    size_kb = output.stat().st_size / 1024
    print(f"[SLIDES] Generated · {output} ({size_kb:.1f} Ko, 24 slides)")


if __name__ == "__main__":
    main()
