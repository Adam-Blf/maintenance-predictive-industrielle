# -*- coding: utf-8 -*-
"""Script 11 · Génération PPTX de soutenance EFREI Mastère DE 2025-26.

Génère `reports/11/presentation.pptx` (12 slides, ~9-10 min).
Structure conforme aux consignes Sarah MALAEB (Bloc 2 RNCP40875).

Usage ::
    python scripts/11_build_pptx.py
"""

from __future__ import annotations

import sys
from pathlib import Path

# Bootstrap avant tout import externe
PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

from src.bootstrap import ensure_dependencies

ensure_dependencies(verbose=False)

# ---------------------------------------------------------------------------
# Imports standard + tiers
# ---------------------------------------------------------------------------
import os
import warnings

warnings.filterwarnings("ignore")

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

try:
    from PIL import Image as PILImage

    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# ---------------------------------------------------------------------------
# Constantes de charte EFREI
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0D, 0x47, 0xA1)
BLUE_VIVID = RGBColor(0x1E, 0x88, 0xE5)
GREEN_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
GRAY_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)
BLACK = RGBColor(0x21, 0x21, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROW_ALT = RGBColor(0xF5, 0xF7, 0xFA)
NAVY_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)

# Dimensions 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Chemins
REPORTS = PROJECT_ROOT / "reports"
ASSETS = PROJECT_ROOT / "assets"
OUT_DIR = REPORTS / "11"

TOTAL_SLIDES = 12

# ---------------------------------------------------------------------------
# Helpers bas niveau
# ---------------------------------------------------------------------------


def _rgb_to_fill(shape, rgb: RGBColor) -> None:
    """Remplit le fond d'une shape avec une couleur unie."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def set_slide_background(slide, rgb: RGBColor) -> None:
    """Applique une couleur de fond à toute la slide via un rectangle couvrant."""
    from pptx.util import Emu

    bg = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0,
        0,
        SLIDE_W,
        SLIDE_H,
    )
    _rgb_to_fill(bg, rgb)
    bg.line.fill.background()
    # Place le rectangle tout en bas de la pile
    sp_tree = slide.shapes._spTree
    sp_tree.remove(bg._element)
    sp_tree.insert(2, bg._element)


def add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = BLACK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = "Calibri",
    wrap: bool = True,
) -> None:
    """Ajoute une zone de texte positionnée en Inches."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def add_title(slide, text: str, color: RGBColor = NAVY, size: int = 24) -> None:
    """Barre de titre bandeau navy en haut (sauf slides 1 et 12).

    Bande 0..0.65" couvrant toute la largeur.
    Titre à l'intérieur (top=0.05", height=0.55"), largeur max 11.5"
    pour laisser la place au logo (left=12.0", width=0.9").
    """
    # Barre colorée pleine largeur
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.65))
    _rgb_to_fill(bar, NAVY)
    bar.line.fill.background()

    # Titre à l'intérieur du bandeau, largeur max 10.7" depuis left=0.15"
    # right edge = 0.15 + 10.7 = 10.85" < 11.0" (debut page number) -> pas d'overlap
    add_text_box(
        slide,
        text,
        left=0.15,
        top=0.05,
        width=10.7,
        height=0.55,
        font_size=size,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.LEFT,
    )


def add_page_number(slide, page_num: int) -> None:
    """Numérotation X / 12 dans le bandeau, à gauche du logo.

    Position : left=11.0" width=0.9" pour rester dans les 11.65" max
    et ne pas chevaucher le logo (left=12.0").
    """
    add_text_box(
        slide,
        f"{page_num} / {TOTAL_SLIDES}",
        left=11.0,
        top=0.07,
        width=0.85,
        height=0.5,
        font_size=11,
        color=RGBColor(0xBD, 0xC3, 0xCB),
        align=PP_ALIGN.RIGHT,
    )


def add_footer(slide) -> None:
    """Footer standard en bas de slide."""
    add_text_box(
        slide,
        "Maintenance Prédictive · EFREI 2025-26 · Adam Beloucif · Emilien Morice",
        left=0.5,
        top=7.05,
        width=12.33,
        height=0.3,
        font_size=9,
        color=GRAY_TEXT,
        align=PP_ALIGN.CENTER,
    )


def add_logo(slide, white: bool = False, cover_slide: bool = False) -> None:
    """Logo EFREI positionné selon le type de slide.

    Slides couverture (cover_slide=True) : logo blanc centré, height=0.6" width=1.2".
    Slides intérieures : logo couleur top-droite, height=0.45" width=0.9",
    left=12.0", top=0.18" (dans la bande 0.65" sans chevaucher le titre).
    Ratio 2:1 strict : width = height * 2.
    """
    if white:
        logo_file = "logo_efrei_white.png"
    else:
        logo_file = "logo_efrei.png"

    logo_path = ASSETS / logo_file
    if not logo_path.exists():
        fallback = ASSETS / "logo_efrei_noir.png"
        if fallback.exists():
            logo_path = fallback
        else:
            return

    try:
        if cover_slide:
            # Slide 1 et 12 : logo blanc centré en haut
            logo_h = Inches(0.6)
            logo_w = Inches(0.6 * 2)  # ratio 2:1
            logo_left = Inches((13.333 - 1.2) / 2)
            logo_top = Inches(0.4)
        else:
            # Slides 2..11 : logo couleur top droite dans le bandeau
            logo_h = Inches(0.45)
            logo_w = Inches(0.45 * 2)  # ratio 2:1 -> width = 0.9"
            logo_left = Inches(12.0)
            logo_top = Inches(0.1)

        slide.shapes.add_picture(
            str(logo_path),
            logo_left,
            logo_top,
            width=logo_w,
            height=logo_h,
        )
    except Exception as e:
        print(f"[logo] warning: {e}")


def add_image(
    slide,
    path: str | Path,
    left: float,
    top: float,
    width: float,
    caption: str | None = None,
    max_height: float = 4.5,
) -> float:
    """Insère une image en Inches. Retourne la hauteur réelle insérée."""
    path = Path(path)
    if not path.exists():
        print(f"[image] WARNING: absent · {path}")
        return 0.0

    # Calcul du ratio si PIL disponible
    target_w = Inches(width)
    if HAS_PIL:
        try:
            with PILImage.open(path) as im:
                iw, ih = im.size
            ratio = ih / iw
            target_h = target_w * ratio
            max_h_emu = Inches(max_height)
            if target_h > max_h_emu:
                target_h = max_h_emu
                target_w = int(target_h / ratio)
        except Exception:
            target_h = None
    else:
        target_h = None

    try:
        if target_h:
            pic = slide.shapes.add_picture(
                str(path), Inches(left), Inches(top), width=target_w, height=target_h
            )
        else:
            pic = slide.shapes.add_picture(
                str(path), Inches(left), Inches(top), width=target_w
            )
    except Exception as e:
        print(f"[image] ERROR inserting {path}: {e}")
        return 0.0

    actual_h = pic.height / 914400  # EMU -> inches

    if caption:
        cap_top = top + actual_h + 0.08
        # Largeur caption = largeur réelle image insérée
        actual_w = pic.width / 914400
        add_text_box(
            slide,
            caption,
            left=left,
            top=cap_top,
            width=actual_w,
            height=0.3,
            font_size=11,
            italic=True,
            color=GRAY_TEXT,
            align=PP_ALIGN.CENTER,
        )

    return actual_h


def add_bullets(
    slide,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    color: RGBColor = BLACK,
    bold_first: bool = False,
) -> None:
    """Ajoute une liste de bullets avec le médiopoint comme puce."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"· {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        run.font.bold = bold_first and i == 0


def add_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left: float,
    top: float,
    width: float,
    height: float,
    highlight_row: int | None = None,
) -> None:
    """Crée un tableau stylé EFREI avec ligne d'en-tête navy."""
    n_rows = len(rows) + 1
    n_cols = len(headers)

    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    col_w = Inches(width) // n_cols
    for col in tbl.columns:
        col.width = col_w

    # En-tête
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.size = Pt(12)
        run.font.name = "Calibri"

    # Lignes de données
    for i, row in enumerate(rows):
        is_highlight = highlight_row is not None and i == highlight_row
        bg = GREEN_LIGHT if is_highlight else (WHITE if i % 2 == 0 else ROW_ALT)
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(11)
            run.font.name = "Calibri"
            run.font.bold = is_highlight
            run.font.color.rgb = BLACK


def _new_slide(prs: Presentation) -> object:
    """Ajoute une slide vierge (layout Blank)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _set_notes(slide, text: str) -> None:
    """Ajoute des notes de présentation sur la slide."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = text
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

# Début du contenu après le bandeau header (0.65")
_CONTENT_TOP = 0.75  # inches - marge sous le bandeau
_FOOTER_TOP = 7.05   # inches - position footer
_MAX_CONTENT_BOTTOM = 6.9  # inches - limite inférieure du contenu (avant footer)


def slide_1_title(prs: Presentation) -> None:
    """Slide 1 · Titre, fond navy plein."""
    slide = _new_slide(prs)
    set_slide_background(slide, NAVY)

    # Logo blanc centré en haut (cover_slide=True -> centré, height=0.6" width=1.2")
    add_logo(slide, white=True, cover_slide=True)

    # Badge RNCP - positionné sous le logo (logo occupe 0.4"..1.0")
    add_text_box(
        slide,
        "RNCP40875 Bloc 2",
        left=0.3,
        top=1.1,
        width=3.5,
        height=0.4,
        font_size=11,
        color=RGBColor(0x90, 0xCA, 0xF9),
        align=PP_ALIGN.LEFT,
    )

    # Titre principal
    add_text_box(
        slide,
        "Système Intelligent Multi-Modèles",
        left=0.5,
        top=1.55,
        width=12.3,
        height=0.85,
        font_size=36,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        "Maintenance Prédictive Industrielle",
        left=0.5,
        top=2.45,
        width=12.3,
        height=0.75,
        font_size=30,
        bold=True,
        color=RGBColor(0x90, 0xCA, 0xF9),
        align=PP_ALIGN.CENTER,
    )

    # Séparateur
    sep = slide.shapes.add_shape(1, Inches(2.5), Inches(3.28), Inches(8.3), Inches(0.04))
    _rgb_to_fill(sep, BLUE_VIVID)
    sep.line.fill.background()

    # Sous-titre
    add_text_box(
        slide,
        "Soutenance Projet Data Science",
        left=0.5,
        top=3.43,
        width=12.3,
        height=0.5,
        font_size=20,
        color=RGBColor(0xBB, 0xDE, 0xFB),
        align=PP_ALIGN.CENTER,
    )

    # Binôme
    add_text_box(
        slide,
        "Adam BELOUCIF (20220055)  ·  Emilien MORICE",
        left=0.5,
        top=4.03,
        width=12.3,
        height=0.45,
        font_size=18,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Infos contexte
    add_text_box(
        slide,
        "M1 Mastère Data Engineering & IA  ·  EFREI  ·  2025-2026",
        left=0.5,
        top=4.58,
        width=12.3,
        height=0.4,
        font_size=15,
        color=RGBColor(0xBB, 0xDE, 0xFB),
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        "Tutrice : Sarah MALAEB",
        left=0.5,
        top=5.08,
        width=12.3,
        height=0.4,
        font_size=14,
        color=RGBColor(0x90, 0xCA, 0xF9),
        align=PP_ALIGN.CENTER,
    )

    _set_notes(
        slide,
        "Bonjour, je suis Adam Beloucif et mon binôme est Emilien Morice. "
        "Nous allons vous présenter notre système intelligent multi-modèles "
        "de maintenance prédictive industrielle, réalisé dans le cadre du "
        "module Projet Data Science du M1 Mastère DE&IA à l'EFREI. "
        "La soutenance dure environ 9-10 minutes.",
    )


def slide_2_context(prs: Presentation) -> None:
    """Slide 2 · Contexte & problématique + pitch."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Contexte & Problématique")
    add_page_number(slide, 2)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Pitch encadré - sous le bandeau (0.65")
    pitch_box = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.72), Inches(12.7), Inches(1.0)
    )
    _rgb_to_fill(pitch_box, NAVY_LIGHT)
    pitch_box.line.color.rgb = BLUE_VIVID

    add_text_box(
        slide,
        'Pitch : "Un système intelligent multi-modèles pour anticiper les pannes à 24h, '
        "qualifier le type de défaillance et estimer la durée de vie restante, "
        'afin d\'aider les équipes maintenance à prioriser leurs interventions."',
        left=0.45,
        top=0.76,
        width=12.4,
        height=0.92,
        font_size=13,
        italic=True,
        color=NAVY,
        align=PP_ALIGN.LEFT,
    )

    # 2 colonnes - début à 1.82" (sous le pitch box)
    # Colonne gauche · enjeux
    add_text_box(
        slide,
        "Enjeux industriels",
        left=0.3,
        top=1.85,
        width=6.0,
        height=0.4,
        font_size=15,
        bold=True,
        color=NAVY,
    )
    add_bullets(
        slide,
        [
            "Coût d'un arrêt machine : 5 000 à 50 000 EUR/h",
            "3 stratégies : corrective, préventive, prédictive",
            "Objectif business : réduire les faux négatifs (pannes non détectées)",
            "ROI : anticiper = planifier = économiser",
        ],
        left=0.3,
        top=2.3,
        width=6.1,
        height=2.3,
        font_size=13,
    )

    # Colonne droite · 3 tâches ML
    add_text_box(
        slide,
        "3 tâches de modélisation",
        left=6.7,
        top=1.85,
        width=6.3,
        height=0.4,
        font_size=15,
        bold=True,
        color=NAVY,
    )

    task_data = [
        ["Tâche", "Cible", "Métrique"],
        ["Classification binaire", "failure_within_24h", "F1 / PR-AUC"],
        ["Classification multiclasse", "failure_type (5 classes)", "Macro-F1"],
        ["Régression RUL", "rul_hours", "MAE / R2"],
    ]
    # top=2.3" + height=1.8" = 4.1" < 7.05" OK
    add_table(
        slide,
        task_data[0],
        task_data[1:],
        left=6.7,
        top=2.3,
        width=6.3,
        height=1.8,
    )

    _set_notes(
        slide,
        "Un arrêt machine non planifié peut coûter jusqu'à 50k EUR/heure. "
        "Nous attaquons le problème sous trois angles : détecter une panne à 24h, "
        "identifier le type de défaillance, estimer le temps de vie restant. "
        "L'enjeu métier principal est de minimiser les faux négatifs.",
    )


def slide_3_user_need(prs: Presentation) -> None:
    """Slide 3 · Besoin utilisateur."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Besoin Utilisateur · Responsable Maintenance")
    add_page_number(slide, 3)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Persona encadré - sous le bandeau
    persona_box = slide.shapes.add_shape(
        1, Inches(0.3), Inches(0.72), Inches(4.5), Inches(0.85)
    )
    _rgb_to_fill(persona_box, NAVY)
    persona_box.line.fill.background()
    add_text_box(
        slide,
        "Persona : Responsable de maintenance d'usine\nObjectif : zéro arrêt imprévu",
        left=0.4,
        top=0.76,
        width=4.3,
        height=0.77,
        font_size=12,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Tableau décisions / indicateurs
    add_text_box(
        slide,
        "Décisions et indicateurs clés",
        left=0.3,
        top=1.7,
        width=6.5,
        height=0.4,
        font_size=15,
        bold=True,
        color=NAVY,
    )
    # top=2.15" + height=2.2" = 4.35" OK
    add_table(
        slide,
        ["Décision métier", "Indicateur fourni"],
        [
            ["Alerte temps réel", "Probabilité panne à 24h + niveau risque"],
            ["Plan d'intervention", "Type de défaillance prédit (5 classes)"],
            ["Priorisation urgences", "Durée de vie restante estimée (heures)"],
            ["Arbitrage budgétaire", "Coût évité : économies estimées en EUR"],
        ],
        left=0.3,
        top=2.15,
        width=8.0,
        height=2.2,
    )

    # KPI tableau de bord
    add_text_box(
        slide,
        "KPI tableau de bord",
        left=8.7,
        top=1.7,
        width=4.4,
        height=0.4,
        font_size=15,
        bold=True,
        color=NAVY,
    )
    add_bullets(
        slide,
        [
            "Machines critiques en temps réel",
            "Plan d'intervention priorisé",
            "Économies réalisées (EUR)",
            "Diagnostic par machine",
            "Détails techniques (jury/DSI)",
        ],
        left=8.7,
        top=2.15,
        width=4.4,
        height=2.2,
        font_size=13,
    )

    # Valeur ajoutée - sous les colonnes
    # top=4.5", height=0.7" => bottom=5.2" < 7.05" OK
    add_text_box(
        slide,
        "Valeur ajoutée : le responsable saisit les valeurs capteurs et reçoit immédiatement "
        "une recommandation actionnelle (aucune action / contrôle 48h / intervention 12-24h).",
        left=0.3,
        top=4.5,
        width=12.7,
        height=0.7,
        font_size=13,
        italic=True,
        color=NAVY,
    )

    _set_notes(
        slide,
        "Notre utilisateur cible est le responsable maintenance. Il a besoin de trois choses : "
        "une alerte avant la panne, un diagnostic précis du type de défaillance, "
        "et une estimation du temps restant pour planifier l'intervention. "
        "Le tableau de bord traduit ces besoins en indicateurs opérationnels directs.",
    )


def slide_4_methodology(prs: Presentation) -> None:
    """Slide 4 · Méthodologie & dataset."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Méthodologie & Dataset")
    add_page_number(slide, 4)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Colonne gauche · méthodologie (début à 0.75")
    add_text_box(
        slide,
        "Approche projet",
        left=0.3,
        top=_CONTENT_TOP,
        width=5.5,
        height=0.4,
        font_size=15,
        bold=True,
        color=NAVY,
    )
    add_bullets(
        slide,
        [
            "Méthode hybride Agile/Kanban (sprints 1 semaine)",
            "Adam : pipeline ML binaire, API, interprétabilité",
            "Emilien : multiclasse, régression, dashboard métier",
            "Outils : GitHub (Adam-Blf), Notion, scripts numérotés",
            "10 scripts Python (02 à 11) + CI/CD GitHub Actions",
        ],
        left=0.3,
        top=1.2,
        width=5.7,
        height=2.3,
        font_size=12,
    )

    # Colonne droite · dataset
    add_text_box(
        slide,
        "Dataset · Kaggle CC0 v3.0",
        left=6.3,
        top=_CONTENT_TOP,
        width=6.7,
        height=0.4,
        font_size=15,
        bold=True,
        color=NAVY,
    )
    add_bullets(
        slide,
        [
            "24 042 lignes x 15 colonnes",
            "7 capteurs numériques + 2 catégorielles + 4 cibles",
            "Déséquilibre : ~14.8% pannes (justifie F1 + PR-AUC)",
            "NaN 2-4% sur 5 capteurs (imputation médiane)",
            "source : tatheerabbas/industrial-machine-predictive-maintenance",
        ],
        left=6.3,
        top=1.2,
        width=6.6,
        height=2.0,
        font_size=12,
    )

    # Tableau colonnes principales - 4 colonnes sur 12.0" depuis left=0.67"
    add_text_box(
        slide,
        "Colonnes principales",
        left=0.67,
        top=3.65,
        width=12.0,
        height=0.4,
        font_size=14,
        bold=True,
        color=NAVY,
    )
    # top=4.1" + height=2.4" = 6.5" < 7.05" OK
    add_table(
        slide,
        ["Colonne", "Type", "Unité", "Rôle"],
        [
            ["vibration_rms", "float", "mm/s", "Feature (top discriminant)"],
            ["temperature_motor", "float", "degC", "Feature (2e discriminant)"],
            ["hours_since_maintenance", "float", "h", "Feature (proxy d'usure)"],
            ["failure_within_24h", "int 0/1", "-", "Cible binaire principale"],
            ["failure_type", "catégorie", "5 classes", "Cible multiclasse"],
            ["rul_hours", "float", "h", "Cible régression (durée vie restante)"],
        ],
        left=0.67,
        top=4.15,
        width=12.0,
        height=2.65,
    )

    _set_notes(
        slide,
        "Nous avons travaillé en binôme avec une méthodologie Agile. "
        "Le dataset Kaggle CC0 contient 24 042 observations et 15 colonnes. "
        "Les deux capteurs les plus discriminants sont la vibration et la température moteur. "
        "Le déséquilibre des classes (~14.8% pannes) justifie notre choix de métriques F1 et PR-AUC.",
    )


def slide_5_eda(prs: Presentation) -> None:
    """Slide 5 · EDA, 2 figures côte à côte."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Analyse Exploratoire (EDA)")
    add_page_number(slide, 5)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Bullets en haut - début à 0.75" (sous bandeau 0.65")
    add_bullets(
        slide,
        [
            "Vibration RMS + température moteur : capteurs les plus discriminants",
            "Mode opératoire 'peak' : taux de panne le plus élevé",
            "5 capteurs avec 2-4% de NaN (imputation médiane robuste aux outliers)",
        ],
        left=0.3,
        top=_CONTENT_TOP,
        width=12.7,
        height=0.85,
        font_size=13,
    )

    # Image gauche : left=0.5", top=1.65", width<=5.8", max_height=4.5"
    # bottom max = 1.65 + 4.5 + 0.08 + 0.3 = 6.53" < 7.05" OK
    heatmap_path = REPORTS / "02" / "eda_correlation_heatmap.png"
    add_image(
        slide,
        heatmap_path,
        left=0.5,
        top=1.65,
        width=5.8,
        caption="Matrice de corrélation capteurs x cible binaire",
        max_height=4.5,
    )

    # Image droite : left=6.83", top=1.65", width<=5.8", max_height=4.5"
    # right edge = 6.83 + 5.8 = 12.63" < 13.333" OK
    scatter_path = REPORTS / "02" / "eda_scatter_vib_temp.png"
    add_image(
        slide,
        scatter_path,
        left=6.83,
        top=1.65,
        width=5.8,
        caption="Scatter vibration x température (zone haute droite = pannes)",
        max_height=4.5,
    )

    _set_notes(
        slide,
        "L'EDA révèle deux capteurs clés : la vibration RMS et la température moteur. "
        "La matrice de corrélation montre des relations modérées, sans redondance critique. "
        "Le scatter plot confirme que les pannes se concentrent dans la zone 'haute droite' "
        "(vibration élevée ET température élevée). Le mode 'peak' est le plus à risque.",
    )


def slide_6_pipeline(prs: Presentation) -> None:
    """Slide 6 · Pipeline IA & preprocessing."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Pipeline IA & Preprocessing")
    add_page_number(slide, 6)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Bullets gauche - début à 0.75"
    add_bullets(
        slide,
        [
            "ColumnTransformer sklearn : anti-leakage garanti (fit train seulement)",
            "Split stratifié 80/20 (stratify=y, random_state=42)",
            "Cross-validation StratifiedKFold 5-fold sur pipeline complet",
            "Déséquilibre : class_weight=balanced (LR/RF) + scale_pos_weight (XGBoost)",
        ],
        left=0.3,
        top=_CONTENT_TOP,
        width=5.5,
        height=2.1,
        font_size=13,
    )

    # Image pipeline droite : left=6.0", width<=7.0" -> right=13.0" OK
    # top=0.75", max_height=5.0" -> bottom max = 0.75+5.0+0.08+0.3=6.13" < 7.05" OK
    pipeline_path = REPORTS / "05" / "diagram_ml_pipeline.png"
    if not pipeline_path.exists():
        pipeline_path = REPORTS / "05" / "diagram_architecture.png"

    add_image(
        slide,
        pipeline_path,
        left=6.0,
        top=_CONTENT_TOP,
        width=7.0,
        caption="Diagramme pipeline ML (preprocessing + train + éval)",
        max_height=5.0,
    )

    # Détail branche preprocessing - sous les bullets
    add_text_box(
        slide,
        "Détail preprocessing",
        left=0.3,
        top=3.0,
        width=5.5,
        height=0.4,
        font_size=14,
        bold=True,
        color=NAVY,
    )
    # top=3.45" + height=1.85" = 5.3" < 7.05" OK
    add_table(
        slide,
        ["Branche", "Étape 1", "Étape 2"],
        [
            ["Numérique (7)", "SimpleImputer(median)", "StandardScaler"],
            ["Catégorielle (2)", "SimpleImputer(mode)", "OneHotEncoder"],
            ["Cibles/ID", "remainder=drop", "(anti-leakage)"],
        ],
        left=0.3,
        top=3.45,
        width=5.5,
        height=1.85,
    )

    _set_notes(
        slide,
        "Le préprocesseur scikit-learn est encapsulé dans un Pipeline avec l'estimateur. "
        "Cela garantit l'absence de data leakage : le fit se fait uniquement sur le train. "
        "La branche numérique impute d'abord puis normalise. "
        "remainder=drop exclut timestamp et identifiants de machine des features.",
    )


def slide_7_models(prs: Presentation) -> None:
    """Slide 7 · Modèles entraînés, tableau comparatif."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Modèles Entraînés · Comparaison")
    add_page_number(slide, 7)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Tableau binaire complet - début à 0.75"
    add_text_box(
        slide,
        "Classification binaire (failure_within_24h)",
        left=0.3,
        top=_CONTENT_TOP,
        width=12.7,
        height=0.4,
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # top=1.2" + height=2.1" = 3.3" OK
    add_table(
        slide,
        ["Modèle", "F1", "ROC-AUC", "PR-AUC", "Latence (ms/sample)"],
        [
            ["Logistic Regression", "0.747", "0.959", "0.838", "0.003"],
            ["Random Forest", "0.863", "0.992", "0.954", "0.021"],
            ["XGBoost (FINAL)", "0.886", "0.995", "0.974", "0.006"],
            ["MLP (scikit-learn)", "0.836", "0.984", "0.909", "0.004"],
        ],
        left=0.3,
        top=1.2,
        width=12.7,
        height=2.1,
        highlight_row=2,  # XGBoost (index 2)
    )

    # MLP note encadrée - top=3.45" + height=0.55" = 4.0" OK
    mlp_box = slide.shapes.add_shape(
        1, Inches(0.3), Inches(3.45), Inches(12.7), Inches(0.55)
    )
    _rgb_to_fill(mlp_box, NAVY_LIGHT)
    mlp_box.line.color.rgb = BLUE_VIVID
    add_text_box(
        slide,
        "MLP : MLPClassifier scikit-learn (64-32-16, relu, adam, early_stopping). "
        "Choix justifié vs TensorFlow/PyTorch : sur-dimensionné pour 24k lignes tabulaires (ADR 0001).",
        left=0.4,
        top=3.48,
        width=12.5,
        height=0.5,
        font_size=12,
        italic=True,
        color=NAVY,
    )

    # Mini-résumé multiclasse + régression
    add_text_box(
        slide,
        "Multiclasse (failure_type · 5 classes)",
        left=0.3,
        top=4.15,
        width=6.0,
        height=0.4,
        font_size=13,
        bold=True,
        color=NAVY,
    )
    # top=4.6" + height=1.0" = 5.6" < 6.7" OK
    add_bullets(
        slide,
        [
            "XGBoost final : Macro-F1 = 0.931, Weighted-F1 = 0.977",
            "Random Forest : Macro-F1 = 0.899",
            "Logistic Regression : Macro-F1 = 0.675",
        ],
        left=0.3,
        top=4.6,
        width=6.1,
        height=1.1,
        font_size=12,
    )

    add_text_box(
        slide,
        "Régression RUL (rul_hours)",
        left=6.8,
        top=4.15,
        width=6.0,
        height=0.4,
        font_size=13,
        bold=True,
        color=NAVY,
    )
    # top=4.6" + height=1.0" = 5.6" < 6.7" OK
    add_bullets(
        slide,
        [
            "Random Forest final : MAE = 9.57h, R2 = 0.673",
            "XGBoost : MAE = 10.72h, R2 = 0.656",
            "Ridge (baseline) : MAE = 20.47h, R2 = 0.140",
        ],
        left=6.8,
        top=4.6,
        width=6.1,
        height=1.1,
        font_size=12,
    )

    _set_notes(
        slide,
        "Nous avons comparé 4 modèles sur chaque tâche. "
        "En binaire, XGBoost remporte le meilleur compromis F1/stabilité CV. "
        "Le MLP scikit-learn suffit pour du tabulaire 24k lignes. "
        "En multiclasse, XGBoost atteint Macro-F1=0.931. "
        "Pour la régression RUL, Random Forest est meilleur avec MAE=9.57h.",
    )


def slide_8_results(prs: Presentation) -> None:
    """Slide 8 · Évaluation comparative, 2 visuels clés."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Évaluation Comparative · Résultats")
    add_page_number(slide, 8)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Bullets synthèse - début à 0.75"
    add_bullets(
        slide,
        [
            "Binaire XGBoost final : F1=0.886, ROC-AUC=0.995, PR-AUC=0.974, CV F1=0.886+/-0.011",
            "Multiclasse XGBoost : Macro-F1=0.931 (bearing 0.918 · hydraulic 0.923 · none 0.987)",
            "Régression RUL Random Forest : R2=0.673, MAE=9.57h, RMSE=15.04h",
        ],
        left=0.3,
        top=_CONTENT_TOP,
        width=12.7,
        height=0.95,
        font_size=13,
    )

    # Image gauche : left=0.5", top=1.78", width<=5.8", max_height=4.5"
    # bottom max = 1.78+4.5+0.08+0.3=6.66" < 7.05" OK
    roc_path = REPORTS / "03" / "roc_curves_comparison.png"
    add_image(
        slide,
        roc_path,
        left=0.5,
        top=1.78,
        width=5.8,
        caption="Courbes ROC superposées (4 modèles binaires)",
        max_height=4.5,
    )

    # Image droite : left=6.83", top=1.78", width<=5.8"
    cm_path = REPORTS / "03" / "confusion_matrix_xgboost.png"
    add_image(
        slide,
        cm_path,
        left=6.83,
        top=1.78,
        width=5.8,
        caption="Matrice de confusion XGBoost (final binaire)",
        max_height=4.5,
    )

    _set_notes(
        slide,
        "La courbe ROC XGBoost est quasiment parfaite (AUC=0.995). "
        "La matrice de confusion montre que le modèle gère bien le déséquilibre : "
        "peu de faux négatifs, ce qui est l'enjeu clé en maintenance. "
        "En multiclasse, chaque type de panne est détecté avec F1 > 0.90.",
    )


def slide_9_interpretability(prs: Presentation) -> None:
    """Slide 9 · Interprétabilité SHAP."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Interprétabilité · SHAP & Feature Importance")
    add_page_number(slide, 9)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Image SHAP gauche : left=0.5", top=0.75", width=6.2", max_height=5.5"
    # Avec caption : bottom max = 0.75+5.5+0.08+0.3=6.63" < 7.05" OK
    shap_path = REPORTS / "04" / "shap_summary_xgboost.png"
    shap_h = add_image(
        slide,
        shap_path,
        left=0.5,
        top=_CONTENT_TOP,
        width=6.2,
        caption="SHAP dot plot XGBoost (contribution par observation)",
        max_height=5.5,
    )

    # Bullets droite - début à 0.75"
    add_text_box(
        slide,
        "Conclusions interprétabilité",
        left=7.0,
        top=_CONTENT_TOP,
        width=5.8,
        height=0.4,
        font_size=15,
        bold=True,
        color=NAVY,
    )
    add_bullets(
        slide,
        [
            "vibration_rms : feature la plus impactante (~30% du gain XGBoost)",
            "temperature_motor : 2e rang (~20%) · confirme l'EDA",
            "Traduction métier : vib > seuil ET temp > seuil = alerte critique",
            "3 niveaux d'interprétabilité : feature importance native + permutation + SHAP",
        ],
        left=7.0,
        top=1.23,
        width=5.8,
        height=2.3,
        font_size=13,
    )

    # Note RNCP encadrée
    note_box = slide.shapes.add_shape(
        1, Inches(7.0), Inches(3.65), Inches(5.8), Inches(0.65)
    )
    _rgb_to_fill(note_box, GREEN_LIGHT)
    note_box.line.color.rgb = GREEN_SUCCESS
    add_text_box(
        slide,
        "Couvre RNCP C3.1/C3.2/C3.3 : explicabilité des modèles, confiance métier, "
        "audit de biais.",
        left=7.1,
        top=3.68,
        width=5.6,
        height=0.59,
        font_size=12,
        color=RGBColor(0x06, 0x5F, 0x46),
    )

    # Image permutation : left=7.0", top=4.45", width=5.8", max_height=1.9"
    # bottom max = 4.45+1.9+0.08+0.3=6.73" < 7.05" OK
    perm_path = REPORTS / "04" / "permutation_importance_xgboost.png"
    add_image(
        slide,
        perm_path,
        left=7.0,
        top=4.45,
        width=5.8,
        caption="Permutation importance XGBoost (perte F1)",
        max_height=1.9,
    )

    _set_notes(
        slide,
        "L'interprétabilité est un critère RNCP. Nous avons trois niveaux : "
        "la feature importance native XGBoost, la permutation importance, et SHAP TreeExplainer. "
        "Les trois s'accordent : vibration et température sont dominants. "
        "Cette convergence renforce la confiance dans le modèle.",
    )


def slide_10_dashboard_api(prs: Presentation) -> None:
    """Slide 10 · Dashboard + API + cost-sensitive."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Dashboard Décisionnel + API · Démo")
    add_page_number(slide, 10)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Colonne gauche · cost-sensitive + figure
    add_text_box(
        slide,
        "Calibration cost-sensitive",
        left=0.3,
        top=_CONTENT_TOP,
        width=5.8,
        height=0.4,
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # Image : left=0.5", top=1.2", width=5.5", max_height=2.8"
    # caption bottom max = 1.2+2.8+0.08+0.3=4.38" -> eco box at 4.5" no overlap
    cost_path = REPORTS / "10" / "cost_threshold_xgboost.png"
    add_image(
        slide,
        cost_path,
        left=0.5,
        top=1.2,
        width=5.5,
        caption="Courbe coût(seuil) XGBoost · minimum à 0.23",
        max_height=2.8,
    )

    # Encadré économies - top=4.5", height=0.9" -> bottom=5.4" < 7.05" OK
    # image caption bottom max = 1.2+2.8+0.08+0.3=4.38" -> eco box at 4.5" => gap 0.12" OK
    eco_box = slide.shapes.add_shape(
        1, Inches(0.3), Inches(4.5), Inches(5.8), Inches(0.9)
    )
    _rgb_to_fill(eco_box, GREEN_LIGHT)
    eco_box.line.color.rgb = GREEN_SUCCESS
    add_text_box(
        slide,
        "Seuil optimal : 0.23 (vs 0.5 défaut)\n"
        "FN=7 x 1000 EUR + FP=256 x 100 EUR = 32 600 EUR\n"
        "Économie : ~12 000 EUR par cycle de scoring",
        left=0.4,
        top=4.53,
        width=5.6,
        height=0.85,
        font_size=12,
        bold=True,
        color=RGBColor(0x06, 0x5F, 0x46),
    )

    # Colonne droite · dashboard onglets
    add_text_box(
        slide,
        "Dashboard Streamlit · 5 onglets",
        left=6.5,
        top=_CONTENT_TOP,
        width=6.5,
        height=0.4,
        font_size=14,
        bold=True,
        color=NAVY,
    )
    # top=1.2" + height=2.3" = 3.5" OK
    add_bullets(
        slide,
        [
            "État du parc : vue temps réel 20 machines, top alertes",
            "Plan d'intervention : table priorisée par urgence",
            "Impact économique : ROI, coûts évités, comparaison avec/sans IA",
            "Diagnostic machine : saisie capteurs -> recommandation temps réel",
            "Détails techniques : EDA, modèles, interprétabilité (sous-onglets)",
        ],
        left=6.5,
        top=1.2,
        width=6.5,
        height=2.3,
        font_size=12,
    )

    # API endpoints
    add_text_box(
        slide,
        "API FastAPI · 4 endpoints",
        left=6.5,
        top=3.65,
        width=6.5,
        height=0.4,
        font_size=14,
        bold=True,
        color=NAVY,
    )
    # top=4.1" + height=1.9" = 6.0" < 7.05" OK
    add_table(
        slide,
        ["Méthode", "Route", "Description"],
        [
            ["GET", "/", "Accueil + liens docs"],
            ["GET", "/health", "Status + model_loaded"],
            ["GET", "/model-info", "Métriques + features"],
            ["POST", "/predict", "Inférence temps réel"],
        ],
        left=6.5,
        top=4.1,
        width=6.5,
        height=1.9,
    )

    # Note backup démo - top=6.45" (bien après fin table API à 6.0") height=0.3" -> bottom=6.75" < 7.05" OK
    add_text_box(
        slide,
        "Backup : screenshots disponibles si démo live indisponible.",
        left=0.5,
        top=6.45,
        width=12.33,
        height=0.3,
        font_size=10,
        italic=True,
        color=GRAY_TEXT,
        align=PP_ALIGN.CENTER,
    )

    _set_notes(
        slide,
        "Le seuil de décision est abaissé à 0.23 : on préfère déclencher des fausses alertes "
        "plutôt que de manquer une vraie panne. Cela économise ~12k EUR par cycle. "
        "Le dashboard Streamlit a 5 onglets orientés métier. "
        "L'API FastAPI expose 4 endpoints avec validation Pydantic automatique.",
    )


def slide_11_limits(prs: Presentation) -> None:
    """Slide 11 · Limites et améliorations."""
    slide = _new_slide(prs)
    set_slide_background(slide, GRAY_BG)
    add_title(slide, "Limites & Pistes d'Amélioration")
    add_page_number(slide, 11)
    add_logo(slide, cover_slide=False)
    add_footer(slide)

    # Colonne gauche · limites - début à 0.75"
    add_text_box(
        slide,
        "Limites actuelles",
        left=0.3,
        top=_CONTENT_TOP,
        width=6.0,
        height=0.4,
        font_size=15,
        bold=True,
        color=RGBColor(0xB9, 0x1C, 0x1C),
    )
    # top=1.2" + height=4.5" = 5.7" < 7.05" OK
    add_bullets(
        slide,
        [
            "Données simulées Kaggle, pas un vrai flux IoT en production",
            "Déséquilibre résiduel malgré class_weight/scale_pos_weight",
            "Optuna documentaire : best_params non réinjectés dans les modèles",
            "MLP scikit-learn : léger overfitting (F1 CV = 0.795 < F1 test = 0.836)",
            "Pas de monitoring de drift (concept drift non détecté)",
            "API CORS permissive (allow_origins=*) - à restreindre en prod",
            "Brier score non persisté (print stdout seulement)",
        ],
        left=0.3,
        top=1.2,
        width=5.9,
        height=4.5,
        font_size=12,
        color=RGBColor(0x4B, 0x05, 0x05),
    )

    # Colonne droite · améliorations
    add_text_box(
        slide,
        "Pistes d'amélioration",
        left=6.7,
        top=_CONTENT_TOP,
        width=6.3,
        height=0.4,
        font_size=15,
        bold=True,
        color=RGBColor(0x06, 0x5F, 0x46),
    )
    # top=1.2" + height=4.5" = 5.7" < 7.05" OK
    add_bullets(
        slide,
        [
            "Réinjection best_params Optuna dans src/models.py",
            "Drift monitoring : Evidently AI, alerte par email",
            "Déploiement cloud : Azure ML / Google Cloud Run",
            "OAuth2 sur l'API + restriction CORS en production",
            "Base de données historique prédictions (PostgreSQL)",
            "Retraining trimestriel automatique (pipeline CI/CD)",
            "Tests d'intégration E2E : valider F1 >= seuil sur test set",
        ],
        left=6.7,
        top=1.2,
        width=6.3,
        height=4.5,
        font_size=12,
        color=RGBColor(0x06, 0x5F, 0x46),
    )

    _set_notes(
        slide,
        "Ces limites sont connues et documentées dans nos ADR. "
        "La principale est que nos données sont simulées, sans vrai drift industriel. "
        "En production, il faudrait ajouter un monitoring Evidently, "
        "un retraining automatique et une gestion OAuth2 sur l'API.",
    )


def slide_12_conclusion(prs: Presentation) -> None:
    """Slide 12 · Conclusion + Q&R, fond navy plein."""
    slide = _new_slide(prs)
    set_slide_background(slide, NAVY)

    # Logo blanc centré en haut (cover_slide=True)
    add_logo(slide, white=True, cover_slide=True)

    # Titre - sous le logo (logo : 0.4"..1.0"), titre commence à 1.1"
    add_text_box(
        slide,
        "Conclusion",
        left=0.5,
        top=1.1,
        width=12.3,
        height=0.6,
        font_size=30,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # Séparateur
    sep = slide.shapes.add_shape(1, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.04))
    _rgb_to_fill(sep, BLUE_VIVID)
    sep.line.fill.background()

    # Bilan projet
    add_text_box(
        slide,
        "Bilan du projet",
        left=0.5,
        top=1.95,
        width=12.3,
        height=0.4,
        font_size=16,
        bold=True,
        color=RGBColor(0x90, 0xCA, 0xF9),
        align=PP_ALIGN.CENTER,
    )

    # top=2.4" + height=2.3" = 4.7" OK
    add_table(
        slide,
        ["Livrable", "Détail"],
        [
            ["12 modèles entraînés", "4 binaire + 4 multiclasse + 4 régression"],
            ["3 tâches couvertes", "Binaire F1=0.886 · Multiclasse Macro-F1=0.931 · RUL R2=0.673"],
            ["Dashboard 5 onglets", "Streamlit orienté métier (responsable maintenance)"],
            ["API 4 endpoints", "FastAPI + Pydantic, latence < 1ms/prédiction"],
            ["RNCP40875 Bloc 2", "C3.1 EDA · C3.2 Prép · C3.3 Éval · C4.3 Empreinte carbone"],
        ],
        left=0.5,
        top=2.4,
        width=12.3,
        height=2.3,
    )

    # Compétences
    add_text_box(
        slide,
        "Compétences acquises",
        left=0.5,
        top=4.85,
        width=12.3,
        height=0.38,
        font_size=14,
        bold=True,
        color=RGBColor(0x90, 0xCA, 0xF9),
        align=PP_ALIGN.CENTER,
    )

    comp_line = (
        "EDA · Preprocessing anti-leakage · Classification binaire & multiclasse · "
        "Régression · DL tabulaire (MLP) · Interprétabilité SHAP · "
        "Dashboard Streamlit · API FastAPI · Gestion de projet Agile"
    )
    add_text_box(
        slide,
        comp_line,
        left=0.5,
        top=5.28,
        width=12.3,
        height=0.5,
        font_size=12,
        color=RGBColor(0xBB, 0xDE, 0xFB),
        align=PP_ALIGN.CENTER,
    )

    # Séparateur 2
    sep2 = slide.shapes.add_shape(1, Inches(3.0), Inches(5.92), Inches(7.3), Inches(0.04))
    _rgb_to_fill(sep2, BLUE_VIVID)
    sep2.line.fill.background()

    # Merci / Q&R
    add_text_box(
        slide,
        "Merci pour votre attention · Questions ?",
        left=0.5,
        top=6.05,
        width=12.3,
        height=0.65,
        font_size=24,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    # IDs étudiants - top=6.78" + height=0.38" = 7.16" < 7.5" OK
    add_text_box(
        slide,
        "Adam BELOUCIF 20220055  ·  Emilien MORICE  ·  M1 DE&IA EFREI 2025-2026",
        left=0.5,
        top=6.78,
        width=12.3,
        height=0.38,
        font_size=11,
        color=RGBColor(0x90, 0xCA, 0xF9),
        align=PP_ALIGN.CENTER,
    )

    _set_notes(
        slide,
        "En conclusion, nous avons construit un système complet : 12 modèles sur 3 tâches, "
        "un dashboard métier et une API de scoring. "
        "Le seuil cost-sensitive de 0.23 économise 12k EUR par cycle. "
        "Ce projet couvre l'intégralité du Bloc 2 RNCP40875. "
        "Nous sommes prêts pour vos questions.",
    )


# ---------------------------------------------------------------------------
# Audit overlaps / out-of-bounds
# ---------------------------------------------------------------------------


def _audit_overlaps(prs: Presentation) -> tuple[int, int]:
    """Détecte overlaps et out-of-bounds. Retourne (n_overlaps, n_out_of_bounds)."""
    sw_emu = prs.slide_width
    sh_emu = prs.slide_height
    n_overlaps = 0
    n_out = 0
    TOL = 9525 * 2  # ~0.02 inch tolerance

    for idx, slide in enumerate(prs.slides, 1):
        rects = []
        for shape in slide.shapes:
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
                if None in (left, top, width, height):
                    continue
                right, bottom = left + width, top + height
                # out of bounds
                if (
                    left < -TOL
                    or top < -TOL
                    or right > sw_emu + TOL
                    or bottom > sh_emu + TOL
                ):
                    n_out += 1
                    print(
                        f"[OOB] slide {idx} shape {shape.shape_type} · "
                        f"L={left/914400:.2f}\" T={top/914400:.2f}\" "
                        f"R={right/914400:.2f}\" B={bottom/914400:.2f}\""
                    )
                rects.append((shape, left, top, right, bottom))
            except Exception:
                continue

        # paires
        for i in range(len(rects)):
            si, l1, t1, r1, b1 = rects[i]
            for j in range(i + 1, len(rects)):
                sj, l2, t2, r2, b2 = rects[j]
                # overlap = intersection rectangle non vide
                overlap_w = min(r1, r2) - max(l1, l2)
                overlap_h = min(b1, b2) - max(t1, t2)
                if overlap_w > TOL and overlap_h > TOL:
                    # whitelist : header bandeau (couvre tout le top) vs titre slide
                    # skip si l'un est exactement le bandeau (largeur >= slide_width et hauteur <= 0.7")
                    is_banner_i = (r1 - l1) >= sw_emu * 0.99 and (b1 - t1) <= 914400 * 0.7
                    is_banner_j = (r2 - l2) >= sw_emu * 0.99 and (b2 - t2) <= 914400 * 0.7
                    if is_banner_i or is_banner_j:
                        continue
                    # whitelist : fond slide (couvre toute la slide)
                    is_bg_i = (r1 - l1) >= sw_emu * 0.99 and (b1 - t1) >= sh_emu * 0.99
                    is_bg_j = (r2 - l2) >= sw_emu * 0.99 and (b2 - t2) >= sh_emu * 0.99
                    if is_bg_i or is_bg_j:
                        continue
                    # whitelist : encadré coloré avec texte posé dessus (AUTO_SHAPE + TEXT_BOX)
                    # Le text box est contenu dans (ou proche de) l'auto shape - pattern intentionnel
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    si_type = si.shape_type
                    sj_type = sj.shape_type
                    AUTO = 1  # MSO_AUTO_SHAPE / RECTANGLE
                    TBOX = 17  # TEXT_BOX
                    if (si_type == AUTO and sj_type == TBOX) or (si_type == TBOX and sj_type == AUTO):
                        # Vérifie que le text box est géométriquement contenu dans l'auto shape
                        if si_type == AUTO:
                            box_l, box_t, box_r, box_b = l1, t1, r1, b1
                            txt_l, txt_t, txt_r, txt_b = l2, t2, r2, b2
                        else:
                            box_l, box_t, box_r, box_b = l2, t2, r2, b2
                            txt_l, txt_t, txt_r, txt_b = l1, t1, r1, b1
                        margin = 914400 * 0.3  # 0.3" tolerance
                        if (
                            txt_l >= box_l - margin
                            and txt_t >= box_t - margin
                            and txt_r <= box_r + margin
                            and txt_b <= box_b + margin
                        ):
                            continue
                    n_overlaps += 1
                    print(
                        f"[OVERLAP] slide {idx} shapes "
                        f"({si.shape_type}, {sj.shape_type}) "
                        f"overlap {overlap_w/914400:.2f}\" x {overlap_h/914400:.2f}\""
                    )

    print(f"\nAudit · {n_overlaps} overlaps detected, {n_out} out of bounds")
    return n_overlaps, n_out


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> None:
    """Point d'entrée principal."""
    print("[11_build_pptx] Démarrage génération PPTX...")

    # Création dossier de sortie
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    # Présentation 16:9
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Génération des 12 slides
    builders = [
        slide_1_title,
        slide_2_context,
        slide_3_user_need,
        slide_4_methodology,
        slide_5_eda,
        slide_6_pipeline,
        slide_7_models,
        slide_8_results,
        slide_9_interpretability,
        slide_10_dashboard_api,
        slide_11_limits,
        slide_12_conclusion,
    ]

    for i, builder in enumerate(builders, start=1):
        print(f"  · Slide {i:2d}/{TOTAL_SLIDES} : {builder.__name__}")
        builder(prs)

    # Audit overlaps avant sauvegarde
    print("\n[audit] Vérification overlaps et out-of-bounds...")
    n_overlaps, n_out = _audit_overlaps(prs)

    # Sauvegarde
    out_path = OUT_DIR / "presentation.pptx"
    prs.save(str(out_path))

    size_mb = out_path.stat().st_size / (1024 * 1024)
    print(f"\n[11_build_pptx] PPTX généré avec succès.")
    print(f"  Chemin : {out_path}")
    print(f"  Slides  : {len(prs.slides)}")
    print(f"  Taille  : {size_mb:.2f} Mo")

    if size_mb < 0.1:
        print("  [WARNING] Fichier très petit : vérifiez que les figures ont bien été insérées.")

    if n_overlaps > 0 or n_out > 0:
        print(f"  [WARNING] Audit non-propre : {n_overlaps} overlaps, {n_out} OOB")


if __name__ == "__main__":
    main()
